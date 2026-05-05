from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── PALETTE ───────────────────────────────────────────────────────────────────
BLU_SCURO   = RGBColor(0x0D, 0x1B, 0x2A)   # sfondo principale
BLU_MEDIO   = RGBColor(0x1B, 0x46, 0x8B)   # accent
BLU_CHIARO  = RGBColor(0x3A, 0x7B, 0xD5)   # highlight
ARANCIO     = RGBColor(0xFF, 0x8C, 0x00)   # contrasto caldo
BIANCO      = RGBColor(0xFF, 0xFF, 0xFF)
GRIGIO_CH   = RGBColor(0xD0, 0xD8, 0xE8)
VERDE       = RGBColor(0x2E, 0xCC, 0x71)
ROSSO_CH    = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completamente vuoto


# ─── HELPER: rettangolo pieno ──────────────────────────────────────────────────
def rect(slide, l, t, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ─── HELPER: textbox ──────────────────────────────────────────────────────────
def tb(slide, text, l, t, w, h,
       size=18, bold=False, color=BIANCO,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


# ─── HELPER: sfondo blu scuro ─────────────────────────────────────────────────
def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, BLU_SCURO)


# ─── HELPER: striscia colorata in cima ────────────────────────────────────────
def top_bar(slide, color=BLU_MEDIO, h=0.08):
    rect(slide, 0, 0, 13.33, h, color)


# ─── HELPER: linea separatrice ────────────────────────────────────────────────
def hline(slide, t, color=BLU_CHIARO, h=0.03):
    rect(slide, 0.5, t, 12.33, h, color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITOLO
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)

# gradiente sinistro
rect(sl, 0, 0, 5, 7.5, BLU_MEDIO)

# logo/istituzione
tb(sl, "Università degli Studi di Salerno",
   0.3, 0.25, 4.4, 0.5, size=11, color=GRIGIO_CH, italic=True,
   align=PP_ALIGN.CENTER)
tb(sl, "Corso di Laurea in Informatica",
   0.3, 0.7, 4.4, 0.5, size=11, color=GRIGIO_CH, italic=True,
   align=PP_ALIGN.CENTER)

hline(sl, 1.3, color=ARANCIO, h=0.05)
rect(sl, 0.3, 1.35, 4.4, 0.05, BLU_CHIARO)

# titolo principale (destra)
tb(sl, "Knowledge Distillation",
   5.5, 1.0, 7.5, 1.0, size=34, bold=True, color=BIANCO,
   align=PP_ALIGN.LEFT)
tb(sl, "vs",
   5.5, 2.0, 7.5, 0.7, size=28, bold=False, color=ARANCIO,
   align=PP_ALIGN.LEFT)
tb(sl, "Reflection Tuning",
   5.5, 2.6, 7.5, 1.0, size=34, bold=True, color=BLU_CHIARO,
   align=PP_ALIGN.LEFT)

tb(sl, "Un confronto sperimentale su modelli linguistici di piccola taglia",
   5.5, 3.7, 7.5, 0.8, size=15, color=GRIGIO_CH, align=PP_ALIGN.LEFT)

hline(sl, 4.6, color=ARANCIO, h=0.04)

tb(sl, "Candidato:  Eliodoro Mascolo",
   5.5, 4.75, 7.5, 0.45, size=13, color=BIANCO)
tb(sl, "Anno Accademico 2024/2025",
   5.5, 5.25, 7.5, 0.45, size=12, color=GRIGIO_CH)

# icone decorative sinistro
tb(sl, "🤖", 0.9, 2.5, 2.5, 1.5, size=72, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)

tb(sl, "Agenda della presentazione",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO,
   align=PP_ALIGN.LEFT)

items = [
    ("01", "Contesto e motivazione",       "Perché ottimizzare i modelli linguistici?"),
    ("02", "Knowledge Distillation",        "Teacher, Student e Dark Knowledge"),
    ("03", "Reflection Tuning",             "Auto-correzione e ragionamento esplicito"),
    ("04", "Setup sperimentale",            "Pipeline, dataset, modelli e metriche"),
    ("05", "Risultati & Analisi",           "LLM-as-a-Judge su 20 domande matematiche"),
    ("06", "Conclusioni & Sviluppi futuri", "Cosa abbiamo imparato"),
]

for i, (num, title, sub) in enumerate(items):
    y = 0.85 + i * 1.05
    rect(sl, 0.4, y, 0.55, 0.78, BLU_CHIARO)
    tb(sl, num, 0.4, y + 0.15, 0.55, 0.5, size=17, bold=True,
       color=BIANCO, align=PP_ALIGN.CENTER)
    tb(sl, title, 1.1, y + 0.05, 5.5, 0.42, size=16, bold=True, color=BIANCO)
    tb(sl, sub,   1.1, y + 0.42, 10.5, 0.38, size=12, color=GRIGIO_CH)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – CONTESTO: IL PROBLEMA DELLA SCALA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "01  |  Contesto & Motivazione",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

tb(sl, "Il paradosso del gigantismo nei Large Language Models",
   0.5, 0.75, 12, 0.5, size=17, bold=True, color=BLU_CHIARO)

# 3 box problema
problems = [
    ("💸", "Costi proibitivi",
     "GPT-4: decine di milioni di $\nper addestramento"),
    ("⚡", "Latenza d'inferenza",
     "Modelli enormi incompatibili\ncon applicazioni real-time"),
    ("🏠", "Edge computing escluso",
     "Smartphone e IoT non possono\nospicare miliardi di parametri"),
]
for i, (ico, title, body) in enumerate(problems):
    x = 0.5 + i * 4.2
    rect(sl, x, 1.4, 3.9, 2.5, BLU_MEDIO)
    tb(sl, ico,   x + 0.15, 1.5,  3.6, 0.8,  size=36, align=PP_ALIGN.CENTER)
    tb(sl, title, x + 0.15, 2.25, 3.6, 0.5,  size=14, bold=True,
       color=ARANCIO, align=PP_ALIGN.CENTER)
    tb(sl, body,  x + 0.15, 2.75, 3.6, 0.95, size=12, color=GRIGIO_CH,
       align=PP_ALIGN.CENTER)

hline(sl, 4.05, color=ARANCIO)
tb(sl, "La domanda chiave:",
   0.5, 4.2, 12, 0.4, size=14, bold=True, color=ARANCIO)
tb(sl,
   "Come ottenere modelli piccoli, veloci e affidabili "
   "senza sacrificare le capacità di ragionamento?",
   0.5, 4.6, 12, 0.7, size=16, color=BIANCO)

tb(sl, "Scaling Laws (Kaplan et al.) → rendimenti marginali decrescenti dopo la soglia ottimale dati/parametri",
   0.5, 5.5, 12, 0.5, size=11, color=GRIGIO_CH, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – KNOWLEDGE DISTILLATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "02  |  Knowledge Distillation",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# schema teacher → student
rect(sl, 0.5, 0.85, 3.2, 1.4, BLU_CHIARO)
tb(sl, "TEACHER\n(Gemini Flash)", 0.5, 0.85, 3.2, 1.4, size=15, bold=True,
   color=BIANCO, align=PP_ALIGN.CENTER)

rect(sl, 5.1, 0.85, 3.2, 1.4, BLU_MEDIO)
tb(sl, "STUDENT\n(Llama 3.2 1B)", 5.1, 0.85, 3.2, 1.4, size=15, bold=True,
   color=BIANCO, align=PP_ALIGN.CENTER)

# freccia
rect(sl, 3.75, 1.35, 1.3, 0.25, ARANCIO)
tb(sl, "▶", 4.5, 1.3, 0.4, 0.4, size=18, color=ARANCIO)
tb(sl, "Dark Knowledge\n(soft labels)", 3.7, 1.65, 1.6, 0.6,
   size=10, color=ARANCIO, align=PP_ALIGN.CENTER)

# punti chiave
kd_points = [
    "Trasferimento di conoscenza da un modello complesso (Teacher) a uno più piccolo (Student)",
    "Il Teacher trasmette distribuzioni di probabilità sulle risposte, non solo l'etichetta corretta",
    "La 'Dark Knowledge' rivela relazioni semantiche e gerarchie di similarità apprese dal Teacher",
    "Response-based distillation: lo Student imita le risposte finali del Teacher (behavior cloning)",
    "Vantaggio: risposte concise, corrette e ben strutturate; adatto a hardware consumer-grade",
    "Limite: lo Student replica anche i pattern di errore del Teacher; nessun meccanismo di auto-correzione",
]
tb(sl, "Principi fondamentali", 0.5, 2.45, 12, 0.4, size=14, bold=True, color=BLU_CHIARO)
for i, pt in enumerate(kd_points):
    y = 2.9 + i * 0.65
    rect(sl, 0.5, y, 0.2, 0.35, ARANCIO)
    tb(sl, pt, 0.85, y - 0.04, 12.0, 0.55, size=12, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – REFLECTION TUNING
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "03  |  Reflection Tuning",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# ciclo iterativo
steps = [
    ("1", "Genera risposta iniziale"),
    ("2", "<thinking> block"),
    ("3", "Auto-valutazione critica"),
    ("4", "Risposta rivista e definitiva"),
]
for i, (n, lbl) in enumerate(steps):
    x = 0.5 + i * 3.1
    rect(sl, x, 0.85, 2.6, 1.1, BLU_MEDIO)
    tb(sl, n,   x + 0.1, 0.88, 0.6, 0.5, size=22, bold=True, color=ARANCIO)
    tb(sl, lbl, x + 0.1, 1.3,  2.3, 0.55, size=12, color=BIANCO)
    if i < 3:
        rect(sl, x + 2.63, 1.25, 0.45, 0.2, ARANCIO)
        tb(sl, "→", x + 2.6, 1.2, 0.5, 0.4, size=16, color=ARANCIO)

rt_points = [
    "Ispirato al Sistema 2 di Kahneman: pensiero lento, deliberato e critico",
    "Il modello produce un blocco <thinking> con ragionamento esplicito prima della risposta",
    "Fine-tuning su dati già annotati con catene di pensiero generate da modelli più grandi",
    "Efficacia dipende fortemente dalla scala: funziona meglio con modelli ≥7B parametri",
    "Vantaggio: auto-correzione su problemi multi-passo; riduzione delle allucinazioni",
    "Limite: a 1B parametri, il modello fatica a mantenere un ragionamento coerente nel blocco <thinking>",
]
tb(sl, "Principi fondamentali", 0.5, 2.15, 12, 0.4, size=14, bold=True, color=BLU_CHIARO)
for i, pt in enumerate(rt_points):
    y = 2.6 + i * 0.65
    rect(sl, 0.5, y, 0.2, 0.35, BLU_CHIARO)
    tb(sl, pt, 0.85, y - 0.04, 12.0, 0.55, size=12, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – CONFRONTO TEORICO
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "Confronto teorico",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

headers = ["Dimensione", "Knowledge Distillation", "Reflection Tuning"]
cols    = [3.0, 4.6, 4.6]
xs      = [0.3, 3.4, 8.1]

# header row
for j, (h, w, x) in enumerate(zip(headers, cols, xs)):
    c = BLU_CHIARO if j > 0 else BLU_SCURO
    rect(sl, x, 0.78, w, 0.55, c)
    tb(sl, h, x + 0.1, 0.82, w - 0.1, 0.45, size=13, bold=True,
       color=BIANCO, align=PP_ALIGN.CENTER)

rows = [
    ("Approccio",        "Imitazione esterna (Teacher)",       "Auto-riflessione interna"),
    ("Conoscenza",       "Trasferita dall'esterno",            "Generata iterativamente"),
    ("Dataset",          "Risposte sintetiche del Teacher",    "Catene di pensiero annotate"),
    ("Requisiti",        "Teacher capace + API",               "Dati con ragionamento esplicito"),
    ("Scala ottimale",   "Efficace anche a 1B param.",         "Richiede ≥7B param."),
    ("Allucinazioni",    "Ereditati dal Teacher",              "Ridotte dall'auto-correzione"),
    ("Velocità infer.",  "Rapida (risposta diretta)",          "Più lenta (blocco thinking)"),
]
for i, row in enumerate(rows):
    y    = 1.4 + i * 0.72
    bg_c = BLU_SCURO if i % 2 == 0 else RGBColor(0x12, 0x28, 0x42)
    for j, (cell, w, x) in enumerate(zip(row, cols, xs)):
        rect(sl, x, y, w, 0.68, bg_c)
        fc = ARANCIO if j == 0 else BIANCO
        tb(sl, cell, x + 0.12, y + 0.1, w - 0.15, 0.55,
           size=12, color=fc, bold=(j == 0))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – SETUP SPERIMENTALE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "04  |  Setup sperimentale",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# pipeline diagram
pipeline = [
    ("OpenThoughts\n114k", BLU_CHIARO),
    ("700 campioni\nextracted", BLU_MEDIO),
    ("Gemini Flash\n(Teacher)", ARANCIO),
    ("QLoRA Fine-tuning\nLlama 3.2 1B", VERDE),
    ("Valutazione\nLLM-as-a-Judge", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (lbl, col) in enumerate(pipeline):
    x = 0.3 + i * 2.5
    rect(sl, x, 0.82, 2.1, 1.15, col)
    tb(sl, lbl, x + 0.05, 0.9, 2.0, 0.95, size=12, bold=True,
       color=BIANCO, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(sl, "→", x + 2.15, 1.18, 0.38, 0.4, size=18, bold=True, color=ARANCIO)

# dettagli
tb(sl, "Dettagli tecnici", 0.5, 2.2, 12, 0.4, size=14, bold=True, color=BLU_CHIARO)

details = [
    ("Modello base", "Llama 3.2 1B-Instruct  (Meta)"),
    ("Training",     "QLoRA — rango r=16, α=32, quantizzazione 4-bit — Google Colab T4"),
    ("Dataset KD",   "700 risposte generate da Gemini Flash (teacher); stesso numero di epoche"),
    ("Dataset RT",   "700 esempi OpenThoughts con blocchi <thinking> da modelli di reasoning"),
    ("Valutazione",  "20 domande matematiche; giudice: Gemini Flash (LLM-as-a-Judge)"),
    ("Metriche",     "Correctness · Reasoning Quality · Clarity  (scala 1–5) + giudizio vittoria"),
]
for i, (k, v) in enumerate(details):
    y = 2.72 + i * 0.68
    rect(sl, 0.5, y, 2.5, 0.55, BLU_CHIARO)
    tb(sl, k, 0.55, y + 0.07, 2.4, 0.42, size=12, bold=True, color=BIANCO)
    tb(sl, v, 3.15, y + 0.07, 9.7, 0.48, size=12, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – RISULTATI QUANTITATIVI
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "05  |  Risultati quantitativi",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# ── tabella punteggi medi ──────────────────────────────────────────
tb(sl, "Punteggi medi (scala 1–5, n=19 domande valide)",
   0.5, 0.78, 8, 0.4, size=13, bold=True, color=BLU_CHIARO)

col_heads = ["Metrica", "Distillation", "Reflection", "Δ (R−D)"]
cw = [2.6, 2.0, 2.0, 1.5]
cx = [0.5, 3.15, 5.2, 7.25]
for j, (h, w, x) in enumerate(zip(col_heads, cw, cx)):
    rect(sl, x, 1.25, w, 0.48, BLU_CHIARO)
    tb(sl, h, x + 0.05, 1.3, w - 0.05, 0.38, size=12, bold=True,
       color=BIANCO, align=PP_ALIGN.CENTER)

score_rows = [
    ("Correctness",       "3.89", "2.74", "−1.16"),
    ("Reasoning Quality", "3.74", "3.05", "−0.68"),
    ("Clarity",           "4.11", "2.42", "−1.68"),
    ("Media complessiva", "3.91", "2.74", "−1.17"),
]
for i, row in enumerate(score_rows):
    y = 1.78 + i * 0.55
    bg_c = RGBColor(0x12, 0x28, 0x42) if i % 2 == 0 else BLU_SCURO
    bold_row = (i == 3)
    for j, (cell, w, x) in enumerate(zip(row, cw, cx)):
        rect(sl, x, y, w, 0.5, bg_c)
        fc = BIANCO
        if j == 3:
            fc = ROSSO_CH
        if j == 1:
            fc = VERDE
        if j == 2:
            fc = ARANCIO
        tb(sl, cell, x + 0.05, y + 0.07, w - 0.05, 0.38,
           size=12, bold=bold_row, color=fc, align=PP_ALIGN.CENTER)

# ── winner distribution ─────────────────────────────────────────
tb(sl, "Giudizi di vittoria (19 domande)",
   9.0, 0.78, 4.0, 0.4, size=13, bold=True, color=BLU_CHIARO)

win_data = [
    ("Distillation", 14, 73.7, VERDE),
    ("Reflection",    5, 26.3, ARANCIO),
    ("Pareggio",      0,  0.0, GRIGIO_CH),
]
# barre orizzontali
bar_x, bar_y0, bar_h, bar_max_w = 9.0, 1.3, 0.52, 4.0
for i, (lbl, cnt, pct, col) in enumerate(win_data):
    y = bar_y0 + i * 0.85
    w = bar_max_w * pct / 100.0
    rect(sl, bar_x, y, max(w, 0.05), bar_h, col)
    tb(sl, f"{lbl}: {cnt} ({pct}%)",
       bar_x + 0.1, y + 0.1, bar_max_w - 0.1, 0.38,
       size=12, bold=True, color=BIANCO)

# big number
tb(sl, "73.7%",  9.0, 3.5, 4.0, 1.1, size=56, bold=True,
   color=VERDE, align=PP_ALIGN.CENTER)
tb(sl, "Distillation",  9.0, 4.55, 4.0, 0.45, size=14,
   color=VERDE, align=PP_ALIGN.CENTER)

tb(sl, "Distillation vince su 14 domande su 19 — il divario maggiore è sulla Clarity (Δ = −1.68)",
   0.5, 5.6, 12.3, 0.5, size=12, color=GRIGIO_CH, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – ANALISI QUALITATIVA
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "05  |  Analisi qualitativa dei risultati",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# colonna sinistra: quando vince KD
rect(sl, 0.4, 0.78, 5.9, 5.8, BLU_MEDIO)
tb(sl, "Quando vince la Distillation  (14/19)",
   0.55, 0.85, 5.6, 0.45, size=14, bold=True, color=VERDE)
kd_cases = [
    "Calcoli diretti e algebra elementare\n(Q4, Q5, Q7, Q8, Q9, Q14, Q17, Q20)",
    "Risposte concise e strutturate grazie\nallo stile del Teacher Gemini Flash",
    "Punteggi massimi su Clarity (media 4.11)\nIndicano output puliti e ben organizzati",
    "Es. Q5 – 5 operai in 8 gg → 10 operai:\nKD: risposta corretta, punteggi 5/5/5\nRT: output incompleto, punteggi 1/1/1",
]
for i, pt in enumerate(kd_cases):
    y = 1.4 + i * 1.15
    rect(sl, 0.55, y, 0.25, 0.4, VERDE)
    tb(sl, pt, 0.95, y - 0.05, 5.1, 0.95, size=11.5, color=BIANCO)

# colonna destra: quando vince RT
rect(sl, 6.9, 0.78, 5.9, 5.8, RGBColor(0x1A, 0x3A, 0x5C))
tb(sl, "Quando vince il Reflection  (5/19)",
   7.05, 0.85, 5.6, 0.45, size=14, bold=True, color=ARANCIO)
rt_cases = [
    "Problemi multi-passo in cui la KD\nsbaglia aritmetica (Q2, Q3)",
    "Il blocco <thinking> funziona da\nscaffolding per passaggi intermedi",
    "Es. Q2 – velocità media del treno:\nKD: errore aritmetico (1/1/2)\nRT: percorre i passi e trova la risposta esatta (5/5/5)",
    "Ragionamento non deterministico\n(Q11, Q12): RT ottiene Reasoning\nQuality 4 vs 1–2 di KD",
]
for i, pt in enumerate(rt_cases):
    y = 1.4 + i * 1.15
    rect(sl, 7.05, y, 0.25, 0.4, ARANCIO)
    tb(sl, pt, 7.45, y - 0.05, 5.1, 0.95, size=11.5, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – DISCUSSIONE & LIMITI
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "05  |  Discussione e limiti",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

tb(sl, "Perché la Distillation vince a questa scala?",
   0.5, 0.78, 12, 0.4, size=15, bold=True, color=BLU_CHIARO)

expl = [
    "Gemini Flash produce risposte elementari corrette e concise → facili da imitare per 1B param.",
    "Le catene di pensiero nei dati RT sono generate da modelli molto più grandi → troppo verbose per 1B",
    "Un modello da 1B non ha abbastanza capacità per mantenere un ragionamento coerente in un lungo blocco <thinking>",
    "DeepSeek-R1, OpenAI o1 e simili funzionano con decine di miliardi di parametri + dataset enormi",
]
for i, e in enumerate(expl):
    rect(sl, 0.5, 1.3 + i * 0.62, 0.25, 0.42, BLU_CHIARO)
    tb(sl, e, 0.9, 1.3 + i * 0.62, 12.0, 0.52, size=13, color=BIANCO)

tb(sl, "Limiti sperimentali da considerare",
   0.5, 3.95, 12, 0.4, size=15, bold=True, color=ARANCIO)

limits = [
    ("Test set ridotto",    "Solo 19 domande valide → tendenze qualitative, non significatività statistica"),
    ("Bias del giudice",    "Gemini Flash è sia Teacher (KD) sia Judge → potenziale stilistic bias"),
    ("Dominio ristretto",   "Solo matematica/logica elementare → risultati potrebbero variare su altri task"),
    ("KD response-based",   "Senza soft targets (inaccessibili via API) → confronto è SFT vs SFT con dati diversi"),
]
for i, (k, v) in enumerate(limits):
    y = 4.48 + i * 0.62
    rect(sl, 0.5, y, 2.4, 0.5, ARANCIO)
    tb(sl, k, 0.55, y + 0.07, 2.3, 0.38, size=11.5, bold=True, color=BLU_SCURO)
    tb(sl, v, 3.05, y + 0.07, 10.0, 0.42, size=12, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – CONCLUSIONI
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "06  |  Conclusioni",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

# box risultato principale
rect(sl, 0.4, 0.82, 12.5, 1.35, BLU_MEDIO)
tb(sl, "Risultato principale",
   0.6, 0.88, 4, 0.4, size=13, bold=True, color=ARANCIO)
tb(sl,
   "A 1B parametri con 700 campioni, la Knowledge Distillation response-based\n"
   "supera nettamente il Reflection Tuning su tutte le metriche (73.7% vittorie).\n"
   "Il Reflection Tuning mostra però un vantaggio su problemi multi-passo dove il\n"
   "blocco <thinking> funge da meccanismo di scaffolding.",
   0.6, 1.25, 12.0, 0.9, size=13, color=BIANCO)

# takeaways
takeaways = [
    ("La scala conta",
     "I vantaggi del Reflection Tuning documentati in letteratura emergono\n"
     "con modelli ≥7B parametri — a 1B quella capacità non è sufficiente"),
    ("Distillation = efficienza per imitazione",
     "Copiare un Teacher capace è una strategia robusta e pratica\n"
     "anche su hardware consumer-grade (Colab T4)"),
    ("Non è una sconfitta del RT",
     "Il risultato non confuta il Reflection Tuning: indica che\n"
     "servono scala e dati significativamente maggiori per sfruttarlo"),
]
for i, (title, body) in enumerate(takeaways):
    y = 2.4 + i * 1.35
    rect(sl, 0.4, y, 0.12, 0.95, BLU_CHIARO)
    tb(sl, title, 0.65, y + 0.02, 12.0, 0.4, size=13, bold=True, color=BLU_CHIARO)
    tb(sl, body,  0.65, y + 0.42, 12.0, 0.75, size=12.5, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – SVILUPPI FUTURI
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
top_bar(sl, BLU_MEDIO, 0.55)
tb(sl, "06  |  Sviluppi futuri",
   0.5, 0.08, 12, 0.45, size=22, bold=True, color=BIANCO)
hline(sl, 0.65)

future = [
    ("🔬", "Scala maggiore",
     "Replicare con modelli da 7B+ parametri per verificare se il vantaggio\n"
     "del Reflection Tuning emerge in modo consistente"),
    ("📊", "Benchmark standardizzati",
     "Ampliare il test set con GSM8K e MATH per valutazioni\n"
     "statisticamente robuste e confrontabili con la letteratura"),
    ("🔗", "Approccio ibrido",
     "Pipeline a 2 stadi: prima KD per la correttezza di base,\n"
     "poi RT per sviluppare capacità di auto-correzione"),
    ("🏠", "Hardware dedicato",
     "Rimuovere il vincolo Colab T4 per training più lungo,\n"
     "dataset più ampio e modelli di scala superiore"),
]
for i, (ico, title, body) in enumerate(future):
    x = 0.5 + (i % 2) * 6.4
    y = 0.85 + (i // 2) * 2.8
    rect(sl, x, y, 5.8, 2.5, BLU_MEDIO)
    tb(sl, ico,   x + 0.2, y + 0.15, 1.0, 0.9, size=36)
    tb(sl, title, x + 1.3, y + 0.2,  4.3, 0.45, size=15, bold=True,
       color=BLU_CHIARO)
    tb(sl, body,  x + 1.3, y + 0.7,  4.3, 1.2, size=12, color=BIANCO)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – SLIDE FINALE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
rect(sl, 0, 0, 13.33, 7.5, BLU_MEDIO)   # sfondo intero

tb(sl, "Grazie per l'attenzione",
   1.5, 1.5, 10.0, 1.2, size=42, bold=True, color=BIANCO,
   align=PP_ALIGN.CENTER)

hline(sl, 3.0, color=ARANCIO, h=0.07)

tb(sl, "Domande?",
   1.5, 3.3, 10.0, 0.8, size=30, color=ARANCIO,
   align=PP_ALIGN.CENTER)

tb(sl, "Eliodoro Mascolo",
   1.5, 4.4, 10.0, 0.55, size=16, color=BIANCO,
   align=PP_ALIGN.CENTER)
tb(sl, "Università degli Studi di Salerno  |  A.A. 2024/2025",
   1.5, 4.9, 10.0, 0.45, size=13, color=GRIGIO_CH,
   align=PP_ALIGN.CENTER)

tb(sl,
   '"Insegnare a un modello piccolo a pensare ad alta voce\n'
   'richiede che quel modello abbia già la capacità di pensare."',
   1.5, 5.6, 10.0, 1.0, size=13, italic=True, color=GRIGIO_CH,
   align=PP_ALIGN.CENTER)

# ─── SALVATAGGIO ──────────────────────────────────────────────────────────────
out = "results/Presentazione_Laurea_Mascolo.pptx"
prs.save(out)
print(f"Salvato: {out}")
